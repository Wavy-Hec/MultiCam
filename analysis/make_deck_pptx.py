#!/usr/bin/env python3
"""Assemble the results deck as PowerPoint — "Part 2" of the July-28 slides,
picking up at that deck's closing header ("Starting the actual method
experiments"). Slide list, captions and stats are reused from
make_deck_html.py at run time so the pptx and the html deck cannot drift.

Run (cvbench env has python-pptx):
  python analysis/make_deck_pptx.py [--out analysis/deck_part2.pptx]
Regenerate after mvueval_slide_stats.py + make_slide_figs.py.
"""
import argparse
import html
import os
import runpy
import sys
import tempfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
FIGS = os.path.join(HERE, "figs_deck")

# run the html builder against a throwaway path and lift its computed globals
_tmp = tempfile.NamedTemporaryFile(suffix=".html", delete=False)
_tmp.close()
_argv = sys.argv
sys.argv = ["make_deck_html.py", _tmp.name]
G = runpy.run_path(os.path.join(HERE, "make_deck_html.py"))
sys.argv = _argv
os.unlink(_tmp.name)
SLIDES, S = G["SLIDES"], G["S"]
TASK_ORDER, PRETTY = G["TASK_ORDER"], G["PRETTY"]

INK = RGBColor(0x14, 0x14, 0x13)
INK2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x89, 0x87, 0x81)
ACCENT = RGBColor(0x2A, 0x78, 0xD6)
NEG = RGBColor(0xC2, 0x3A, 0x39)
OK = RGBColor(0x1D, 0x6B, 0x1D)

W, H = Inches(13.333), Inches(7.5)


def text_box(slide, x, y, w, h, runs, size=14, align=PP_ALIGN.LEFT,
             space_after=6):
    """runs: str, or list of paragraphs; a paragraph is str or [(text, opts)]."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if isinstance(para, str):
            para = [(para, {})]
        for txt, opts in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(opts.get("size", size))
            r.font.bold = opts.get("bold", False)
            r.font.italic = opts.get("italic", False)
            r.font.color.rgb = opts.get("color", INK)
    return box


def new_slide(prs, title=None, sub=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if title:
        text_box(s, Inches(0.55), Inches(0.22), W - Inches(1.1), Inches(0.7),
                 [[(title, {"size": 26, "bold": True})]])
    if sub:
        text_box(s, Inches(0.55), Inches(0.92), W - Inches(1.1), Inches(0.5),
                 [[(sub, {"size": 12.5, "color": INK2})]])
    return s


def add_figure(slide, path, top=Inches(1.45), max_h=Inches(4.85),
               max_w=Inches(12.2)):
    pic = slide.shapes.add_picture(path, 0, top, height=max_h)
    if pic.width > max_w:
        scale = max_w / pic.width
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
    pic.left = int((W - pic.width) / 2)
    return pic


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=os.path.join(HERE, "deck_part2.pptx"))
    args = ap.parse_args()

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ---- title -------------------------------------------------------------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    text_box(s, Inches(1.2), Inches(2.2), W - Inches(2.4), Inches(1.6),
             [[("Multi-camera Video Understanding", {"size": 40, "bold": True})],
              [("Part 2 — the method experiments", {"size": 24, "color": INK2})]],
             align=PP_ALIGN.CENTER)
    text_box(s, Inches(1.2), Inches(4.6), W - Inches(2.4), Inches(1.6),
             [[("Picks up at the July-28 deck's closing header "
                "(“Starting the actual method experiments”).",
                {"size": 14, "color": INK2})],
              [("All numbers in this deck: MVU-Eval full pool (1,824 q), "
                "InternVL3-8B, this repo's harness, 4 passes, question-level "
                "permutation tests. Levels are not comparable to Part 1 "
                "(different harness config) — shapes replicate.",
                {"size": 12.5, "color": MUTED})]],
             align=PP_ALIGN.CENTER)

    # ---- corrections to Part 1 --------------------------------------------
    sv, ct = S["single_view"], S["blind"]["contingency"]["overall"]
    s = new_slide(prs, "Two Part-1 numbers, corrected",
                  "Both are now computed by checked-in scripts; quote these "
                  "versions.")
    text_box(
        s, Inches(0.8), Inches(1.7), W - Inches(1.6), Inches(4.6),
        [[("“Best single view +14 to +33”  →  "
           f"+{sv['best'] - sv['luck_best_of_k']:.1f} over luck",
           {"size": 18, "bold": True, "color": NEG})],
         [(f"Best view {sv['best']:.1f} vs sequential "
           f"{sv['sequential_all_views']:.1f} is +21.7 raw, but picking the "
           f"best of K views wins {sv['luck_best_of_k']:.1f} by luck alone. "
           "The honest view-selector headroom is the difference. The old "
           "figures predate the stats bug fix and never subtracted luck.",
           {"size": 14, "color": INK2})],
         [("", {})],
         [("“50–65% wrong both ways”  →  "
           f"{ct['wrong_both']:.1f}% under the documented convention",
           {"size": 18, "bold": True, "color": NEG})],
         [(f"Probabilistic pass-pair convention (written down in "
           f"slide_stats.json): wrong-both {ct['wrong_both']:.1f}%, images "
           f"hurt {ct['vision_hurt']:.1f}%, images helped "
           f"{ct['vision_helped']:.1f}%, right regardless "
           f"{ct['right_both']:.1f}%. The 50–65% figure only appears "
           "under a strict all-4-passes-wrong convention nobody had written "
           "down.", {"size": 14, "color": INK2})]],
        space_after=10)

    # ---- the 14 result slides, same order as the html deck ----------------
    for fn, title, cap in SLIDES:
        path = fn if os.path.isabs(fn) else os.path.join(FIGS, fn)
        s = new_slide(prs, html.unescape(title))
        add_figure(s, path)
        text_box(s, Inches(0.7), Inches(6.45), W - Inches(1.4), Inches(0.95),
                 [[(html.unescape(cap), {"size": 12.5, "color": INK2})]])

    # ---- ours vs the paper, as a real table -------------------------------
    P = S["paper"]
    ours8 = S["arrangement"]["mvu_full"]["cvbench_native"]
    pp = S["parity32pv"]["cvbench_native"]
    s = new_slide(prs, f"Ours vs the MVU-Eval paper ({P['source']})",
                  "InternVL3-8B. Paper: 32 frames/video, ≤720px, vLLM. "
                  "Ours: 448px tiles, HF chat, 4 passes at temperature 0.7 "
                  "(decoding not matched).")
    rows = [("Accuracy, percent", "Paper", "Ours, 8/video",
             "Ours, 32/video (parity)")]
    rows.append(("Overall", f"{P['internvl3_8b']['overall']:.1f}",
                 f"{ours8['acc']:.1f}", f"{pp['acc']:.1f}"))
    for t in TASK_ORDER:
        rows.append((PRETTY[t], f"{P['internvl3_8b']['by_task'][t]:.1f}",
                     f"{ours8['by_task'][t]['acc']:.1f}",
                     f"{pp['by_task'][t]['acc']:.1f}"))
    rows.append(("Random / human", f"{P['random']:.1f} / {P['human']:.1f}",
                 f"{S['arrangement']['chance']['mvu_full']:.1f} / —",
                 "—"))
    tbl = s.shapes.add_table(len(rows), 4, Inches(1.6), Inches(1.6),
                             Inches(10.1), Inches(4.7)).table
    for j, wd in enumerate((3.4, 2.2, 2.25, 2.25)):
        tbl.columns[j].width = Inches(wd)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12 if i else 12.5)
            para.font.bold = i in (0, 1)
            para.font.color.rgb = INK if i else INK2
    text_box(s, Inches(0.7), Inches(6.5), W - Inches(1.4), Inches(0.9),
             [[("We beat the paper's number at every budget including their "
                "own — harness quality, not frames, separates the runs. "
                f"Open item: the paper puts Qwen2.5-VL-7B at "
                f"{P['qwen2_5_vl_7b']['overall']:.1f}, ABOVE its InternVL "
                "number; every ordering we measure is the reverse.",
                {"size": 12.5, "color": INK2})]])

    # ---- what to build next -----------------------------------------------
    svt = S["single_view"]["by_task"]
    s = new_slide(prs, "What to build next — ranked by measured headroom")
    text_box(
        s, Inches(0.8), Inches(1.5), W - Inches(1.6), Inches(5.3),
        [[("1. View selection", {"size": 17, "bold": True, "color": OK}),
          (f"  — the only measured headroom left: "
           f"+{sv['best'] - sv['luck_best_of_k']:.1f} honest points, "
           f"concentrated in KIR (best view {svt['MVU-KIR']['best']:.1f}), "
           f"RAG ({svt['MVU-RAG']['best']:.1f}) and Comparison.",
           {"size": 14, "color": INK2})],
         [("2. A per-clip floor for the CLIP sampler", {"size": 17, "bold": True, "color": OK}),
          ("  — on mvu_178 it kept ZERO frames of the gold clip and went "
           "0/4 (slide 16). Evidence deletion is its one failure mode; the fix "
           "is one line in the top-k allocation.", {"size": 14, "color": INK2})],
         [("3. A Qwen2.5-VL-7B run in this harness", {"size": 17, "bold": True, "color": OK}),
          ("  — settles the ordering flip vs the paper; backend is one "
           "HF download + an uncomment away.", {"size": 14, "color": INK2})],
         [("", {})],
         [("Dead ends, measured:", {"size": 17, "bold": True, "color": NEG}),
          ("  more frames (flat-to-negative across a 4.7× budget range, "
           "incl. exact paper parity) and stitching variants (flat at three "
           "budgets, the one significant per-task delta reverses).",
           {"size": 14, "color": INK2})]],
        space_after=14)

    prs.save(args.out)
    print(f"wrote {args.out} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()

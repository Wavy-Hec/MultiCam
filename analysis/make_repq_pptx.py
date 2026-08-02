"""Render the representative-questions doc as a PowerPoint deck (importable
into Google Slides), in the same visual grammar as the July-28 reference deck:
one message per slide — big question-style title, grey context line, one
exhibit, and a bottom takeaway with the load-bearing numbers in bold.

Reads the SAME sources as make_repq_doc.py (manifest, pools, slide_stats,
sampler_diag) so the two deliverables cannot drift apart.

Usage:
  python3 analysis/make_repq_pptx.py [--out analysis/repq_deck.pptx]
"""
import argparse
import io
import json
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)

from pptx import Presentation                                   # noqa: E402
from pptx.util import Inches, Pt, Emu                           # noqa: E402
from pptx.dml.color import RGBColor                             # noqa: E402
from pptx.enum.text import PP_ALIGN                             # noqa: E402

import make_repq_doc as R                                       # noqa: E402

INK = RGBColor(0x14, 0x14, 0x13)
GREY = RGBColor(0x6B, 0x6A, 0x66)
ACCENT = RGBColor(0x2A, 0x78, 0xD6)
GOLD_BG = RGBColor(0x1D, 0x6B, 0x1D)
W, H = Inches(13.333), Inches(7.5)


def add_text(slide, x, y, w, h, runs, size=12, color=INK, bold=False,
             align=PP_ALIGN.LEFT, line_spacing=1.15):
    """runs: str, or list of paragraphs, each a list of (text, bold, color)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    if isinstance(runs, str):
        runs = [[(runs, bold, color)]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if isinstance(para, str):
            para = [(para, bold, color)]
        for text, b, c in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = b
            r.font.color.rgb = c
            r.font.name = "Arial"
    return box


def new_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    add_text(slide, Inches(0.55), Inches(0.32), W - Inches(1.1), Inches(0.7),
             title, size=26, bold=True)
    if subtitle:
        add_text(slide, Inches(0.55), Inches(1.02), W - Inches(1.1),
                 Inches(0.6), subtitle, size=11.5, color=GREY)
    return slide


def takeaway(slide, runs):
    add_text(slide, Inches(0.55), Inches(6.75), W - Inches(1.1), Inches(0.6),
             [runs] if runs and isinstance(runs[0], tuple) else runs, size=13)


def pil_to_stream(img, q=80):
    buf = io.BytesIO()
    img.convert("RGB").save(buf, "JPEG", quality=q)
    buf.seek(0)
    return buf


def thumb_grid(slide, rec, video, x, y, w, max_thumbs=6):
    media = (R.vids_of(rec) if video else R.imgs_of(rec))[:max_thumbs]
    root = R.MVU_MEDIA if video else R.AAB_MEDIA
    cols = 2 if len(media) > 1 else 1
    cell_w = (w - Inches(0.15) * (cols - 1)) / cols
    cx, cy, row_h = x, y, Emu(0)
    for i, m in enumerate(media, 1):
        try:
            img = (R.video_middle_frame if video else R.image_thumb)(
                os.path.join(root, m), 240)
        except Exception:
            continue
        ih = Emu(int(cell_w * img.height / img.width))
        pic = slide.shapes.add_picture(pil_to_stream(img), cx, cy,
                                       width=cell_w)
        add_text(slide, cx, cy + pic.height, cell_w, Inches(0.22),
                 f"{'Video' if video else 'View'} {i}", size=8, color=GREY)
        row_h = max(row_h, pic.height + Inches(0.24))
        if i % cols == 0:
            cx, cy = x, cy + row_h
            row_h = Emu(0)
        else:
            cx += cell_w + Inches(0.15)
        if cy > H - Inches(1.6):
            break
    n_more = len((R.vids_of(rec) if video else R.imgs_of(rec))) - len(media)
    if n_more > 0:
        add_text(slide, x, min(cy + row_h, H - Inches(1.2)), w, Inches(0.25),
                 f"+{n_more} more (see doc / media bundle)", size=9,
                 color=GREY)


def question_body(slide, rec, tag, chips_line, video, x, w):
    """One flowing textbox — paragraphs can never overlap each other."""
    qt = rec["question"]
    if len(qt) > 380:
        qt = qt[:380] + " … (full text in the doc)"
    gold = str(rec["answer"]).strip().upper()
    paras = [[(tag + f"  ·  id {rec['id']}", True, ACCENT)],
             [("", False, INK)],
             [('"' + qt + '"', False, INK)],
             [("", False, INK)]]
    for o in rec["options"][:8]:
        is_gold = o.strip().upper().startswith(gold + ".")
        o_show = o if len(o) < 110 else o[:110] + "…"
        paras.append([(o_show + ("   ← gold" if is_gold else ""),
                       is_gold, GOLD_BG if is_gold else INK)])
    if len(rec["options"]) > 8:
        paras.append([(f"(+{len(rec['options']) - 8} more options)",
                       False, GREY)])
    if chips_line:
        paras.append([("", False, INK)])
        paras.append([("correct passes /4:  ", True, GREY),
                      (chips_line, False, GREY)])
    add_text(slide, x, Inches(1.55), w, Inches(5.0), paras, size=11.5,
             line_spacing=1.2)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=os.path.join(HERE, "repq_deck.pptx"))
    args = ap.parse_args()

    S = json.load(open(R.STATS))
    diag = json.load(open(R.DIAG))
    gate = R.gate_stats()
    M = json.load(open(R.MANIFEST_DEFAULT)) if hasattr(R, "MANIFEST_DEFAULT") \
        else json.load(open(os.path.join(HERE, "repq_doc_manifest.json")))
    mvu = {r["id"]: r for r in json.load(open(R.MVU_POOL))}
    aab = {r["id"]: r for r in json.load(open(R.AAB_POOL))}
    ego_o2i = {r["orig_id"]: r["id"]
               for r in json.load(open(R.AAB_EGO_POOL))}

    picks_mvu, picks_aab = {}, {}
    for p in M["mvu_picks"]:
        picks_mvu.setdefault(p["task_type"], []).append(p)
    for p in M["aab_picks"]:
        picks_aab.setdefault(p["task_type"], []).append(p)

    picked_mvu = {p["id"] for p in M["mvu_picks"]}
    picked_aab = {ego_o2i[mvu_or_aab["orig_id"]]
                  for mvu_or_aab in (aab[p["id"]] for p in M["aab_picks"])
                  if mvu_or_aab.get("source") == "EgoHumans"}
    chips = R.load_chips(picked_mvu, picked_aab)

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 1 — title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(s, Inches(1.2), Inches(2.6), W - Inches(2.4), Inches(1.4),
             "Representative questions by category", size=40, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.2), Inches(4.0), W - Inches(2.4), Inches(0.9),
             [[("MVU-Eval (1,824 q, 8 categories) · All-Angles (2,132 q, "
                "6 categories) · InternVL3-8B, 4 passes", False, GREY)],
              [("questions picked by deterministic, results-blind rules — "
                "no cherry-picking possible", False, GREY)]],
             size=14, align=PP_ALIGN.CENTER)

    # 2 — how chosen
    s = new_slide(prs, "How were the questions chosen?",
                  "The selection never sees model results and has no "
                  "randomness — every pick re-derivable from the rules.")
    add_text(s, Inches(0.55), Inches(1.7), W - Inches(1.1), Inches(4.5), [
        [("Pick primitive: ", True, INK),
         ("sort a stratum by (question length, id), take the middle element; "
          "ties break to the lowest id.", False, INK)],
        [("Slot A — modal exemplar: ", True, INK),
         ("the median-length question at the category's most common video "
          "count and source mix.", False, INK)],
        [("Slot B — defect exemplar: ", True, INK),
         ("when ≥20% of a category carries a text flag (options naming "
          "videos, duplicated options, >6 options), the median-length "
          "flagged record.", False, INK)],
        [("Slot C — sync exemplar: ", True, INK),
         ("where the category intersects the 128 synchronized nuScenes "
          "6-camera questions.", False, INK)],
        [("", False, INK)],
        [("Per-question results are printed on every card for transparency, "
          "but selection never consulted them.", False, GREY)],
    ], size=15, line_spacing=1.4)

    # 3 — MVU at a glance
    d = S["defects"]
    arr = S["arrangement"]["mvu_full"]
    s = new_slide(prs, "What does MVU-Eval actually contain?",
                  "1,824 questions · 2–13 clips each · chance "
                  f"{S['arrangement']['chance']['mvu_full']}% · majority "
                  f"floor {S['arrangement']['majority_floor']['mvu_full']['acc']:.1f}%")
    add_text(s, Inches(0.55), Inches(1.8), W - Inches(1.1), Inches(4.0), [
        [(f"{d['mvu_names_video']['names_a_video']:,} questions (70%) ",
          True, ACCENT),
         ("name a specific video in text or options — \"which video\" is "
          "part of the task.", False, INK)],
        [(f"{S['meta']['sync_n']} questions (7%) ", True, ACCENT),
         ("are true synchronized 6-camera rigs (nuScenes); the other 1,696 "
          "are non-synchronized clip sets (275 of those are segments of one "
          "source video).", False, INK)],
        [("Gold letters are skewed: ", False, INK),
         (f"A {d['mvu_gold_letter_dist'].get('A')} · "
          f"B {d['mvu_gold_letter_dist'].get('B')} · "
          f"C {d['mvu_gold_letter_dist'].get('C')} · "
          f"D {d['mvu_gold_letter_dist'].get('D')}", True, INK)],
        [("", False, INK)],
        [("Overall by arm:  ", True, INK),
         (f"blind {S['blind']['mvu_full']['acc']:.1f} · "
          f"sequential {arr['cvbench_native']['acc']:.1f} · "
          f"stitching {arr['centralized']['acc']:.1f} · "
          f"decentralized {arr['per_stream']['acc']:.1f} · "
          f"frame_select {S['clip']['frame_select']['acc']:.1f} · "
          f"clip_top1 {S['clip']['clip_select_top1']['acc']:.1f} · "
          f"temporal {S['clip']['temporal_weighted']['acc']:.1f}",
          False, INK)],
    ], size=15, line_spacing=1.5)
    takeaway(s, [("Only 7% of MVU-Eval is synchronized multi-camera; for "
                  "most of the rest the task is cross-video search.",
                  True, INK)])

    # 4..11 — MVU categories
    mvu_cats = sorted(picks_mvu, key=lambda c: -len([r for r in mvu.values()
                                                     if r["task_type"] == c]))
    for cat in mvu_cats:
        cat_recs = [r for r in mvu.values() if r["task_type"] == cat]
        comp = R.composition(cat_recs)
        pt = S["per_task"].get(cat, {})
        ia = S["blind"]["images_add_by_task"].get(cat, {})
        bt = lambda blk: blk["by_task"].get(cat, {}).get("acc")
        f = comp["flag_counts"]
        sub = (f"{comp['n']} questions · videos/q "
               + " ".join(f"{k}:{v}" for k, v in comp["k_dist"].items())
               + f" · sources {R.fam_str(comp['families'])}"
               + f" · names-a-video {100 * f['names_video'] / comp['n']:.0f}%"
               + f" · chance {comp['chance']}%")
        s = new_slide(prs, f"{R.MVU_NAMES.get(cat, cat)}  ({cat})", sub)
        p0 = picks_mvu[cat][0]
        rec = mvu[p0["id"]]
        question_body(s, rec, p0["tag"],
                      R.chip_text(chips["mvu"].get(rec["id"], {}),
                                  R.MVU_METHODS), True,
                      Inches(0.55), Inches(6.6))
        thumb_grid(s, rec, True, Inches(7.5), Inches(1.55), Inches(5.3))
        seq, cen = bt(arr["cvbench_native"]), bt(arr["centralized"])
        bl = S["blind"]["mvu_full"]["by_task"].get(cat, {}).get("acc")
        takeaway(s, [
            (f"sequential {seq:.1f} · stitching Δ {pt.get('delta', 0):+.2f} "
             f"(p={pt.get('p', 1):.2f}) · blind {bl:.1f} — images add "
             f"{ia.get('delta', 0):+.1f} (p={ia.get('p', 1):.3f})",
             True, INK),
            (f"   ({len(picks_mvu[cat])} example(s) in the doc)",
             False, GREY)])

    # 12 — AAB at a glance
    ego_n = sum(1 for r in aab.values() if r.get("source") == "EgoHumans")
    aarr = S["arrangement"]["aab170"]
    s = new_slide(prs, "What does All-Angles actually contain?",
                  f"2,132 questions · still images, 2–5 synchronized cameras "
                  f"· chance {S['arrangement']['chance']['aab170']}% · "
                  f"majority floor "
                  f"{S['arrangement']['majority_floor']['aab170']['acc']:.1f}%")
    add_text(s, Inches(0.55), Inches(1.8), W - Inches(1.1), Inches(4.0), [
        [(f"{2132 - ego_n:,} of 2,132 questions (92%) ", True, ACCENT),
         ("come from Ego-Exo4D, whose images are blocked on a pending "
          "license — every accuracy below is on the 170-question EgoHumans "
          "slice.", False, INK)],
        [("", False, INK)],
        [("By arm (EgoHumans-170):  ", True, INK),
         (f"blind {S['blind']['aab170']['acc']:.1f} · "
          f"sequential {aarr['cvbench_native']['acc']:.1f} · "
          f"stitching {aarr['centralized']['acc']:.1f} · "
          f"decentralized {aarr['per_stream']['acc']:.1f}", False, INK)],
    ], size=15, line_spacing=1.5)
    takeaway(s, [("All six categories are represented in the EgoHumans "
                  "slice (14–40 q each); full-pool text exemplars are in "
                  "the doc.", True, INK)])

    # 13..18 — AAB categories
    aab_cats = sorted(picks_aab, key=lambda c: -len([r for r in aab.values()
                                                     if r["task_type"] == c]))
    for cat in aab_cats:
        cat_recs = [r for r in aab.values() if r["task_type"] == cat]
        comp = R.composition(cat_recs, video=False)
        bt = lambda blk: blk["by_task"].get(cat, {}).get("acc")
        ego_pick = next((p for p in picks_aab[cat]
                         if aab[p["id"]].get("source") == "EgoHumans"), None)
        sub = (f"{comp['n']} questions in the full pool · views/q "
               + " ".join(f"{k}:{v}" for k, v in comp["k_dist"].items())
               + f" · chance {comp['chance']}%")
        s = new_slide(prs, cat.replace("AAB-", "").replace("-", " ")
                      + f"  ({cat})", sub)
        if ego_pick:
            rec = aab[ego_pick["id"]]
            question_body(s, rec, ego_pick["tag"],
                          R.chip_text(chips["aab"].get(
                              ego_o2i.get(rec["orig_id"], -1), {}),
                              R.AAB_METHODS), False,
                          Inches(0.55), Inches(6.6))
            thumb_grid(s, rec, False, Inches(7.5), Inches(1.55), Inches(5.3))
        bl = S["blind"]["aab170"]["by_task"].get(cat, {}).get("acc")
        takeaway(s, [
            (f"blind {R.fmt_pct(bl)} · sequential "
             f"{R.fmt_pct(bt(aarr['cvbench_native']))} · stitching "
             f"{R.fmt_pct(bt(aarr['centralized']))} · decentralized "
             f"{R.fmt_pct(bt(aarr['per_stream']))}  (EgoHumans slice)",
             True, INK)])

    # 19 — sampler answers
    fs = diag["frame_select"]
    g, gt = fs["gold_named"], diag["clip_select_top1"]["gold_named"]
    s = new_slide(prs, "The CLIP frame sampler — encoder & thresholds",
                  "Answers to the three questions asked in the thread.")
    add_text(s, Inches(0.55), Inches(1.7), W - Inches(1.1), Inches(4.6), [
        [("Encoder:  ", True, INK),
         ("openai/clip-vit-base-patch32 (HF AutoModel) — identical on all "
          f"{fs['n_rows']:,} full-pool rows.", False, INK)],
        [("Feature-distance thresholds:  ", True, INK),
         ("none. Not distance-based keyframing (that is the separate MOG2 "
          "motion sampler). Ours: 32 uniform candidates per clip → one CLIP "
          "text-image cosine pass against the raw question → keep the global "
          "top-64. No per-clip floor — a clip can get 0 frames.", False, INK)],
        [("", False, INK)],
        [("Does it align with relevance?  ", True, INK),
         (f"On the {g['n']} questions whose gold option names one video: "
          f"frame_select keeps the relevant clip {g['recall_kept']:.0f}% of "
          f"the time but concentrates only mildly ({g['mean_share']:.3f} "
          f"share vs {g['mean_uniform_share']:.3f} uniform); accuracy "
          f"{g['acc_kept']:.1f}% kept vs {g['acc_dropped']:.1f}% dropped. "
          f"clip_top1 picks the right clip {gt['recall_kept']:.0f}% (random "
          f"{gt['random_pick_recall']:.0f}%): {gt['acc_kept']:.1f}% when "
          f"right vs {gt['acc_dropped']:.1f}% when wrong.", False, INK)],
    ], size=15, line_spacing=1.45)
    takeaway(s, [("The selector, not the harness, is the bottleneck — "
                  f"scorer gate recall@1 {gate['CLIP B/32']['recall1']:.1f}% "
                  f"(CLIP) vs {gate['CLIP B/32']['random_floor']:.1f}% "
                  "random.", True, INK)])

    # 20..23 — the four contact sheets
    qe = M["sampler_ids"]
    captions = {
        "aligned": "QE-aligned — (tied) highest concentration: CLIP piles "
                   "frames onto the gold-named clip.",
        "typical": "QE-typical — the median case: barely departs from "
                   "uniform.",
        "missed": "QE-missed — the gold-named clip got ZERO frames: the "
                  "evidence the question depends on was silently deleted.",
        "degenerate": "QE-degenerate — at K=2 the candidate pool can never "
                      "exceed the budget, so selection is a no-op.",
    }
    from PIL import Image
    for kind in ("aligned", "typical", "missed", "degenerate"):
        qid = qe[kind]
        s = new_slide(prs, f"Qualitative: {captions[kind].split(' — ')[0]}",
                      captions[kind].split(" — ", 1)[1] + f"  (question id "
                      f"{qid}; frames the model actually received)")
        img = Image.open(os.path.join(R.RES, "figs_repq_sampler",
                                      f"mvu_{qid}_frameselect.png"))
        max_w, max_h = W - Inches(1.1), Inches(5.1)
        scale = min(max_w / Emu(Inches(img.width / 96)),
                    max_h / Emu(Inches(img.height / 96)), 1.0)
        s.shapes.add_picture(
            pil_to_stream(img, 85), Inches(0.55), Inches(1.5),
            width=Emu(int(Inches(img.width / 96) * scale)))

    # 24 — where the videos live
    s = new_slide(prs, "Watching the actual videos",
                  "Every clip and image behind every card, byte-identical.")
    add_text(s, Inches(0.55), Inches(1.8), W - Inches(1.1), Inches(4.0), [
        [("Bundle on the cluster:  ", True, INK),
         ("CVBench/bench/results/repq_media/ — one folder per question "
          "(mvu_<id> / aab_<id>), files prefixed with their Video/View "
          "number; index.json maps folders to question text.", False, INK)],
        [("Fetch:  ", True, INK),
         ("scp -r <user>@<cluster>:~/CVBench/bench/"
          "results/repq_media ~/Downloads/", False, INK)],
        [("", False, INK)],
        [("Full doc (all 33 cards, thumbnails, per-pass results): ",
          True, INK), ("repq_doc.html / repq_doc.md, same numbers as this "
                       "deck.", False, INK)],
    ], size=15, line_spacing=1.5)

    prs.save(args.out)
    print(f"wrote {args.out} ({os.path.getsize(args.out) / 1e6:.1f} MB, "
          f"{len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()

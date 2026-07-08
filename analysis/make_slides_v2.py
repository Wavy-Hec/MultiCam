#!/usr/bin/env python3
"""Build the REFRESHED multi-camera VLM progress deck (2026-06-22).

Supersedes make_slides.py / multicam_progress.pptx. Reflects everything since:
  - the clean Qwen CrossView re-run (8192 budget + abstain parser; the old 0/20
    Event-Ordering "result" was a scoring bug, now resolved),
  - the 1,033-Q MEVA scale-up (both models),
  - the failure-mechanism taxonomy read off the traces,
  - the blind-vs-video FLIP (video helps on CVBench, HURTS on CrossView),
  - the input/output pipeline (how video reaches the model, how answers are scored),
  - EgoExo4D now downloaded locally.

Figures are regenerated from the verified result JSONs into analysis/slides_v2_figs/
so the deck always matches the data. CVBench overall accuracies are verified
constants (see ORIENTATION.md §8). 16:9, ~5 bullets/slide, jargon glossed.

Two modes from one source (no number drift):
  full deck  ->  analysis/multicam_progress_v2.pptx       (16 slides)
  --concise  ->  analysis/multicam_progress_short.pptx    (8 slides, the essentials)

Run:  ~/anaconda3/envs/cvbench/bin/python analysis/make_slides_v2.py [--concise]
"""
import json
import os
import sys
from collections import defaultdict

CONCISE = "--concise" in sys.argv

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
A = os.path.join(REPO, "analysis")
ER = os.path.join(REPO, "Video-R1", "src", "r1-v", "eval_results")
FIGS = os.path.join(A, "slides_v2_figs")
OUT = os.path.join(A, "multicam_progress_short.pptx" if CONCISE else "multicam_progress_v2.pptx")
os.makedirs(FIGS, exist_ok=True)

# palette
NAVY_HEX, RED_HEX, GREEN_HEX, GRAY_HEX = "#1F3A5F", "#C0392B", "#1E7E34", "#555555"
QWEN_C, IVL_C = "#1F77B4", "#E67E22"   # qwen blue, internvl orange

# ---------------------------------------------------------------- data helpers
def load(path):
    d = json.load(open(path))
    return d["results"] if isinstance(d, dict) and "results" in d else (
        list(d.values()) if isinstance(d, dict) else d)


def acc(recs):
    n = len(recs)
    return (100.0 * sum(1 for r in recs if r.get("correct")) / n) if n else 0.0


def by_task(recs):
    b = defaultdict(list)
    for r in recs:
        t = (r.get("task_type") or "").replace("CrossView-MEVA-", "")
        if t:
            b[t].append(r)
    return {t: acc(rs) for t, rs in b.items()}


def by_orig_cams(recs):
    b = defaultdict(list)
    for r in recs:
        c = r.get("orig_num_cameras")
        if c is not None:
            b[c].append(r)
    return {c: (acc(rs), len(rs)) for c, rs in sorted(b.items())}


# ---- load the verified result JSONs (CrossView; computed from data) ----
Q1033 = load(os.path.join(ER, "eval_crossview_meva1033_qwen3vl.json"))
IV1033_raw = load(os.path.join(A, "crossview_internvl3_meva1033_normalized.json"))
Q60 = load(os.path.join(ER, "eval_crossview_subset_qwen3vl.json"))

# attach task_type + orig_num_cameras to InternVL-1033 via id join with Qwen
meta = {str(r["id"]): r for r in Q1033}
IV1033 = []
for r in IV1033_raw:
    m = meta.get(str(r["id"]))
    if m:
        r = dict(r, task_type=m.get("task_type"), orig_num_cameras=m.get("orig_num_cameras"))
    IV1033.append(r)

# failure-mechanism counts (recompute the tags used in curate_failure_examples.py)
import re
ROMAN = ["I", "II", "III", "IV", "V", "VI"]


def _asc_letter(options):
    for opt in options or []:
        seq = [s for s in re.findall(r"\b(I{1,3}|IV|V?I{0,3})\b", opt.split(".", 1)[-1]) if s]
        if seq and seq == ROMAN[:len(seq)]:
            return opt.strip()[0]
    return None


def _has_ans(r):
    return "<answer>" in (r.get("output") or "") and bool(str(r.get("prediction") or "").strip())


def _mech(r):
    if not _has_ans(r):
        return "truncation /\nno answer"
    if (r.get("task_type") or "").endswith("Event-Ordering"):
        if _asc_letter(r.get("options")) and str(r.get("prediction") or "").strip()[:1].upper() == _asc_letter(r.get("options")):
            return "given-order\nbias"
    t = (r.get("think") or r.get("output") or "").lower()
    if (r.get("num_videos") or 0) >= 2 and ("the video" in t or "video 1" in t) and not {v for v in ("video 2", "video 3", "video 4") if v in t}:
        return "single-video\nshortcut"
    return "wrong\nreasoning"


MECH = defaultdict(int)
EO_TRUNC = 0
for r in Q1033:
    if not r.get("correct"):
        m = _mech(r)
        MECH[m] += 1
        if m.startswith("truncation") and (r.get("task_type") or "").endswith("Event-Ordering"):
            EO_TRUNC += 1

# ---------------------------------------------------------------- figures
def _style(ax):
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.tick_params(length=0)
    ax.set_axisbelow(True)
    ax.grid(axis="y", color="#dddddd", lw=0.8)


def _bars_labels(ax, bars, fmt="{:.0f}%", dy=1.0, size=11):
    for b in bars:
        ax.annotate(fmt.format(b.get_height()), (b.get_x() + b.get_width() / 2, b.get_height() + dy),
                    ha="center", va="bottom", fontsize=size, color=GRAY_HEX, fontweight="bold")


def fig_headline():
    groups = ["CVBench\n45-Q (assoc.)", "CrossView\n60-Q", "CrossView\n1,033-Q (MEVA)"]
    qwen = [62.2, acc(Q60), acc(Q1033)]
    ivl = [55.6, 38.3, acc(IV1033)]
    x = range(len(groups)); w = 0.38
    fig, ax = plt.subplots(figsize=(8.6, 4.5))
    b1 = ax.bar([i - w / 2 for i in x], qwen, w, label="Qwen3-VL-8B-Thinking", color=QWEN_C)
    b2 = ax.bar([i + w / 2 for i in x], ivl, w, label="InternVL3-8B", color=IVL_C)
    _bars_labels(ax, b1); _bars_labels(ax, b2)
    ax.axhline(28, ls="--", lw=1.2, color=GRAY_HEX); ax.text(2.46, 29, "chance ~28%", fontsize=9, color=GRAY_HEX)
    ax.set_xticks(list(x)); ax.set_xticklabels(groups, fontsize=11)
    ax.set_ylabel("accuracy (%)"); ax.set_ylim(0, 75)
    ax.set_title("Both models: strong on CVBench, near-floor on multi-camera CrossView", fontsize=12, color=NAVY_HEX)
    ax.legend(frameon=False, fontsize=10, loc="upper right")
    _style(ax); fig.tight_layout(); p = os.path.join(FIGS, "headline.png"); fig.savefig(p, dpi=150); plt.close(fig)
    return p


def fig_blind_flip():
    groups = ["CVBench 45-Q", "CrossView 60-Q"]
    video = [62.2, acc(Q60)]
    blind = [40.0, acc(load(os.path.join(ER, "eval_crossview_subset_qwen3vl_novideo.json")))]
    x = range(len(groups)); w = 0.38
    fig, ax = plt.subplots(figsize=(8.0, 4.5))
    b1 = ax.bar([i - w / 2 for i in x], video, w, label="with video (8 frames/clip)", color=NAVY_HEX)
    b2 = ax.bar([i + w / 2 for i in x], blind, w, label="blind (text only, no video)", color="#9aa7b4")
    _bars_labels(ax, b1); _bars_labels(ax, b2)
    ax.annotate("video +22 pts", (0, max(video[0], blind[0]) + 6), ha="center", color=GREEN_HEX, fontweight="bold", fontsize=11)
    ax.annotate("video -5 pts (!)", (1, max(video[1], blind[1]) + 6), ha="center", color=RED_HEX, fontweight="bold", fontsize=11)
    ax.set_xticks(list(x)); ax.set_xticklabels(groups, fontsize=11)
    ax.set_ylabel("Qwen3-VL accuracy (%)"); ax.set_ylim(0, 75)
    ax.set_title("The flip: video HELPS on CVBench but HURTS on CrossView", fontsize=12, color=NAVY_HEX)
    ax.legend(frameon=False, fontsize=10, loc="upper right")
    _style(ax); fig.tight_layout(); p = os.path.join(FIGS, "blind_flip.png"); fig.savefig(p, dpi=150); plt.close(fig)
    return p


def fig_eo_fix():
    fig, ax = plt.subplots(figsize=(7.2, 4.4))
    bars = ax.bar(["old scoring\n(truncation + stray-letter\nparser bug)", "after fix\n(8192 budget +\nabstain parser)"],
                  [0.0, 35.0], color=[RED_HEX, GREEN_HEX], width=0.55)
    _bars_labels(ax, bars, dy=0.8)
    ax.set_ylabel("Qwen Event-Ordering accuracy (%)"); ax.set_ylim(0, 45)
    ax.set_title("Resolved: Qwen's 0/20 Event-Ordering was a scoring bug, not a result", fontsize=11.5, color=NAVY_HEX)
    ax.annotate("0/20 was mechanically\nguaranteed (gold is never 'A')", (0, 4), ha="center", fontsize=9, color=GRAY_HEX)
    ax.annotate("7/20 — really orders events", (1, 37.5), ha="center", fontsize=9, color=GRAY_HEX)
    _style(ax); fig.tight_layout(); p = os.path.join(FIGS, "eo_fix.png"); fig.savefig(p, dpi=150); plt.close(fig)
    return p


def fig_task_flip():
    order = ["Event-Ordering", "Spatial", "Temporal"]
    qt, it = by_task(Q1033), by_task(IV1033)
    qv = [qt.get(t, 0) for t in order]; iv = [it.get(t, 0) for t in order]
    x = range(len(order)); w = 0.38
    fig, ax = plt.subplots(figsize=(8.4, 4.5))
    b1 = ax.bar([i - w / 2 for i in x], qv, w, label="Qwen3-VL", color=QWEN_C)
    b2 = ax.bar([i + w / 2 for i in x], iv, w, label="InternVL3", color=IVL_C)
    _bars_labels(ax, b1); _bars_labels(ax, b2)
    ax.set_xticks(list(x)); ax.set_xticklabels(order, fontsize=11)
    ax.set_ylabel("accuracy (%)  ·  1,033-Q MEVA"); ax.set_ylim(0, 65)
    ax.set_title("The temporal flip: the two models have opposite strengths", fontsize=12, color=NAVY_HEX)
    ax.legend(frameon=False, fontsize=10)
    ax.annotate("Temporal: InternVL 51% vs Qwen 25%", (2, 57), ha="center", color=RED_HEX, fontsize=9.5, fontweight="bold")
    _style(ax); fig.tight_layout(); p = os.path.join(FIGS, "task_flip.png"); fig.savefig(p, dpi=150); plt.close(fig)
    return p


def fig_mechanisms():
    items = sorted(MECH.items(), key=lambda kv: -kv[1])
    labels = [k for k, _ in items]; vals = [v for _, v in items]
    fig, ax = plt.subplots(figsize=(8.6, 4.5))
    bars = ax.barh(range(len(labels)), vals, color=[NAVY_HEX, RED_HEX, "#8e7cc3", "#76a5af"][:len(labels)])
    ax.set_yticks(range(len(labels))); ax.set_yticklabels(labels, fontsize=10.5)
    ax.invert_yaxis()
    for i, v in enumerate(vals):
        ax.annotate(str(v), (v + 4, i), va="center", fontsize=11, fontweight="bold", color=GRAY_HEX)
    ax.set_xlabel(f"# of Qwen failures (of {sum(vals)} wrong on the 1,033-Q pool)")
    ax.set_xlim(0, max(vals) * 1.15)
    ax.set_title("Most failures are genuine wrong reasoning — not scoring/budget artifacts", fontsize=11.5, color=NAVY_HEX)
    ax.annotate(f"{EO_TRUNC} of these are Event-Ordering\n(long chains still exceed 8192 tokens)",
                (max(vals) * 0.42, 1), fontsize=9, color=GRAY_HEX)
    ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False); ax.tick_params(length=0)
    fig.tight_layout(); p = os.path.join(FIGS, "mechanisms.png"); fig.savefig(p, dpi=150); plt.close(fig)
    return p


def fig_camera_curve():
    qc, ic = by_orig_cams(Q1033), by_orig_cams(IV1033)
    cams = sorted(set(qc) | set(ic))
    fig, ax = plt.subplots(figsize=(8.6, 4.5))
    ax.plot(cams, [qc.get(c, (None,))[0] for c in cams], "-o", color=QWEN_C, label="Qwen3-VL")
    ax.plot(cams, [ic.get(c, (None,))[0] for c in cams], "-s", color=IVL_C, label="InternVL3")
    for c in cams:
        n = qc.get(c, (0, 0))[1]
        ax.annotate(f"n={n}", (c, 2), ha="center", fontsize=7.5, color=GRAY_HEX)
    ax.axhline(28, ls="--", lw=1, color=GRAY_HEX); ax.text(cams[-1], 29, "chance", fontsize=8, color=GRAY_HEX, ha="right")
    ax.set_xlabel("original # of synchronized cameras (before the <=4 cap)")
    ax.set_ylabel("accuracy (%)  ·  1,033-Q MEVA"); ax.set_ylim(0, 80)
    ax.set_title("Multi-camera is hard: accuracy stays near the floor as cameras grow", fontsize=12, color=NAVY_HEX)
    ax.legend(frameon=False, fontsize=10)
    _style(ax); fig.tight_layout(); p = os.path.join(FIGS, "camera_curve.png"); fig.savefig(p, dpi=150); plt.close(fig)
    return p


FIG_HEADLINE = fig_headline()
FIG_BLIND = fig_blind_flip()
FIG_EO = fig_eo_fix()
FIG_TASKFLIP = fig_task_flip()
FIG_MECH = fig_mechanisms()
FIG_CAMS = fig_camera_curve()
print("figures:", *(os.path.basename(p) for p in (FIG_HEADLINE, FIG_BLIND, FIG_EO, FIG_TASKFLIP, FIG_MECH, FIG_CAMS)))

# ---------------------------------------------------------------- pptx helpers
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1E, 0x7E, 0x34)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _add_par(tf, text, size=18, bold=False, color=GRAY, level=0, first=False, space_after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size, run.font.bold, run.font.color.rgb = Pt(size), bold, color
    return p


def _title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.32), SW - Inches(1.2), Inches(1.0))
    _add_par(tb.text_frame, text, size=27, bold=True, color=NAVY, first=True)


def _fill_bullets(tf, items, body_size):
    first = True
    for it in items:
        text, level = it if isinstance(it, tuple) else (it, 0)
        bold = text.startswith("**")
        text = text.replace("**", "")
        color = ACCENT if text.startswith("⚠") else (GREEN if text.startswith("✓") else GRAY)
        prefix = "" if (bold and level == 0) else ("•  " if level == 0 else "–  ")
        _add_par(tf, prefix + text, size=body_size - 2 * level, bold=bold,
                 color=color, level=level, first=first, space_after=8)
        first = False


def bullets_slide(title, items, body_size=18):
    s = prs.slides.add_slide(BLANK)
    _title(s, title)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.5), SW - Inches(1.4), SH - Inches(2.0))
    tb.text_frame.word_wrap = True
    _fill_bullets(tb.text_frame, items, body_size)
    return s


def place_image(slide, path, left, top, max_w, max_h):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w, h = max_w, max_w / ar
    if h > max_h:
        h, w = max_h, max_h * ar
    slide.shapes.add_picture(path, Emu(int(left + (max_w - w) / 2)),
                             Emu(int(top + (max_h - h) / 2)), width=Emu(int(w)), height=Emu(int(h)))


def caption(slide, text, color=GRAY, top=None):
    top = top if top is not None else SH - Inches(1.0)
    tb = slide.shapes.add_textbox(Inches(0.7), top, SW - Inches(1.4), Inches(0.95))
    tb.text_frame.word_wrap = True
    _add_par(tb.text_frame, text, size=14, color=color, first=True)


def image_slide(title, img, cap, cap_color=GRAY):
    s = prs.slides.add_slide(BLANK)
    _title(s, title)
    place_image(s, img, Inches(0.7), Inches(1.4), SW - Inches(1.4), Inches(4.6))
    caption(s, cap, color=cap_color)
    return s


def image_then_bullets(title, img, bullets, body_size=15, img_h=Inches(3.85), img_w=None):
    s = prs.slides.add_slide(BLANK)
    _title(s, title)
    iw = img_w or (SW - Inches(1.4))
    place_image(s, img, Inches(0.7), Inches(1.35), iw, img_h)
    btop = Inches(1.35) + img_h + Inches(0.12)
    tb = s.shapes.add_textbox(Inches(0.7), btop, SW - Inches(1.4), SH - btop - Inches(0.2))
    tb.text_frame.word_wrap = True
    _fill_bullets(tb.text_frame, bullets, body_size)
    return s


def image_left_bullets_right(title, img, bullets, body_size=16):
    s = prs.slides.add_slide(BLANK)
    _title(s, title)
    place_image(s, img, Inches(0.6), Inches(1.5), Inches(7.0), Inches(5.2))
    tb = s.shapes.add_textbox(Inches(7.9), Inches(1.6), SW - Inches(8.5), SH - Inches(2.0))
    tb.text_frame.word_wrap = True
    _fill_bullets(tb.text_frame, bullets, body_size)
    return s


def title_slide(title, sub_lines):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.3), SW - Inches(1.6), Inches(2.4))
    tb.text_frame.word_wrap = True
    _add_par(tb.text_frame, title, size=36, bold=True, color=NAVY, first=True)
    for ln in sub_lines:
        _add_par(tb.text_frame, ln, size=19, color=GRAY, space_after=4)
    return s


# ----------------------------------------------------------------- slides
# Each slide is a no-arg builder that appends itself to `prs`. Two ordered
# registries (FULL / CONCISE) select which run, so both decks share one source.

def s_title():
    title_slide("Multi-Camera VLM Benchmarking — Progress Update",
                ["CVBench  +  CrossView (UT Austin multi-camera dataset)",
                 "Now at 1,033 questions · clean re-run · failure taxonomy · EgoExo4D added",
                 "Hector Lugo  ·  for Dr. Datta & Harsh  ·  2026-06-22"])


def s_what_full():
    bullets_slide("What we're testing", [
        "Two \"thinking\" vision-language models — they write out their reasoning before answering, so we can read WHY they fail",
        ("Qwen3-VL-8B-Thinking   and   InternVL3-8B", 1),
        "A \"blind\" control: the SAME questions with the video removed (text only) — isolates the language prior",
        "Two benchmarks side by side: CVBench (unrelated clips) and CrossView (same scene, many synchronized cameras)",
        "**Goal: not just a score — an interpretable picture of HOW the models fail, read from their reasoning traces**",
    ])


def s_what_concise():
    bullets_slide("What we're testing", [
        "Two \"thinking\" VLMs (**Qwen3-VL-8B-Thinking**, **InternVL3-8B**) that write their reasoning before answering — so we read WHY they fail",
        "**CVBench** — 1–4 UNRELATED clips (45-Q)   ·   **CrossView** — the SAME scene from many synchronized cameras (now **1,033 Q**, from UT Austin MEVA)",
        "A **blind** (text-only) control isolates the language prior",
        "**Goal: an interpretable picture of multi-camera failure, not just a score**",
    ], body_size=18)


def s_benchmarks():
    bullets_slide("Two benchmarks, one harness", [
        "**CVBench** — questions over 1–4 UNRELATED clips   (45-Q subset; \"can it connect separate clips?\")",
        "**CrossView** — the SAME scene from many synchronized cameras, from UT Austin MEVA   (now a 1,033-Q pool; \"can it fuse simultaneous views?\")",
        "Both run through the identical thinking harness → directly comparable",
        "Three question types on CrossView: **Event-Ordering, Spatial, Temporal**",
        "Scale-up this round: CrossView 60-Q → **1,033 questions** (both models), so per-type / per-camera numbers are now solid",
    ])


def s_inputs():
    bullets_slide("How the video reaches the model (INPUTS)", [
        "Each clip is reduced to **8 still frames**, then the question is appended — the model never sees full motion video",
        "**Qwen leg** (eval_thinking.py): frames sampled by even spacing INCLUDING first & last; each clip preceded by a text label \"Video k:\"",
        "**InternVL leg** (lmms-eval): frames sampled at segment MIDPOINTS (skips first & last); each clip wrapped by a small marker image (\"This is video k\")",
        "⚠ The two legs don't even sample the same frames or use the same markers — a real asymmetry we track, not yet equalized",
        "Vision-token budget ≈ 11.5k (Qwen) vs 10.3k (InternVL) for a 4-video question",
        "Detail: analysis/input_pipeline.md · analysis/ORIENTATION.md §4",
    ], body_size=16)


def s_outputs():
    bullets_slide("How answers are scored (OUTPUTS)", [
        "Models must emit  <think> … </think>  then  <answer> X </answer>  — we read both",
        "**The fix this round:** if a trace runs out of budget and never writes <answer>, the scorer now ABSTAINS instead of grabbing a stray letter",
        ("the old parser grabbed the first \"a/b/c/d\" in the prose → fabricated answers on truncated traces (see next slide)", 1),
        "Every question now logs token accounting: text tokens, video tokens (count of <|video_pad|>), total, # frames",
        "Two parsers give 30.9% (eval-time) vs 32.0% (re-parse) on the 1,033-Q pool — a ~1pp, 12-question recovery gap, documented",
        "Detail: analysis/ORIENTATION.md §5",
    ], body_size=16)


def s_headline():
    image_then_bullets("Headline numbers", FIG_HEADLINE, [
        "**Both models score well on CVBench (60s%) but fall to the floor (~30s%) on multi-camera CrossView**",
        "CrossView 1,033-Q: Qwen 30.9% (32.0% re-parse) · InternVL3 34.0% — barely above chance (~28%)",
        "Fusing simultaneous camera views of one scene is the hard, unsolved skill this benchmark isolates",
    ], body_size=15)


def s_flip():
    image_then_bullets("The blind control reveals a flip", FIG_BLIND, [
        "**On CVBench, video adds +22 pts** — the answer genuinely needs the frames",
        "**On CrossView, video HURTS Qwen** (blind 36.7% > video 31.7%) — the extra views add confusion the model can't resolve, and the text options alone are sometimes more guessable",
        "This is the central tension: more cameras should help, but current models do worse with them",
    ], body_size=15)


def s_eofix():
    image_then_bullets("Resolved: Qwen's \"0/20\" was a scoring bug, not a result", FIG_EO, [
        "Old run: every Event-Ordering trace ran past the 2048-token limit, never finished, and the scorer defaulted to \"A\" — but gold is never \"A\", so 0/20 was mechanically guaranteed",
        "**Fixed:** raised the budget to 8192 + abstain-on-truncation parser → 54/60 traces now finish; Event-Ordering is a real 7/20",
        "The clean 60-Q overall (31.7%) matches the old number only by coincidence — the composition is now legitimate. Old file archived as *.contaminated-*.bak",
    ], body_size=14)


def s_cameras():
    image_slide("CrossView: accuracy stays near the floor as cameras grow",
                FIG_CAMS,
                "1,033-Q MEVA, both models, vs the TRUE synchronized-camera count (before the <=4 cap). "
                "No recovery as cameras increase — direct evidence that fusing many simultaneous views is the hard core. "
                "(High-camera bins are small-n, labeled.)")


def s_taskflip():
    image_then_bullets("By question type — and the temporal flip", FIG_TASKFLIP, [
        "**The two models have nearly opposite strengths by task type**",
        "Temporal: InternVL3 51% vs Qwen 25%;  Event-Ordering & Spatial: Qwen ahead",
        "Both confuse video-level timestamps with event-level time → temporal grounding is a shared weak spot",
    ], body_size=15)


def s_mech_full():
    image_then_bullets("Failure taxonomy at scale (read from the traces)", FIG_MECH, [
        "**Most failures are genuine wrong reasoning (626) — the model sees the frames and concludes wrongly**, not scoring/budget artifacts",
        "Truncation/no-answer (69) is now an honest abstain, concentrated in long Event-Ordering chains (60)",
        "given-order bias (13) and single-video shortcut (6) round out the tail",
        "Digest: analysis/crossview_meva1033_out/failure_examples_curated.md",
    ], body_size=15)


def s_mech_concise():
    image_then_bullets("Why they fail (read from the traces)", FIG_MECH, [
        "**Most failures are genuine wrong reasoning (626)** — the model sees the frames but concludes wrongly; not scoring/budget artifacts",
        "Tail modes: long Event-Ordering chains truncate (60), given-order bias (13), single-clip shortcutting (6)",
        "Digest: analysis/crossview_meva1033_out/failure_examples_curated.md",
    ], body_size=15)


def s_case():
    s = prs.slides.add_slide(BLANK)
    _title(s, "Worked failure case: counting across 4 cameras")
    fe = os.path.join(A, "failure_examples")
    cells = [("Video 1 — store display", "id57_video1_frames.jpg"),
             ("Video 2 — store display", "id57_video2_frames.jpg"),
             ("Video 3 — store display", "id57_video3_frames.jpg"),
             ("Video 4 — drinks INSIDE the cart (missed)", "id57_video4_frames.jpg")]
    gx, gy = Inches(0.7), Inches(1.45)
    cw, ch = (SW - Inches(1.4)) / 2, Inches(2.0)
    for i, (lab, fn) in enumerate(cells):
        r, c = divmod(i, 2)
        left, top = gx + c * cw, gy + r * (ch + Inches(0.32))
        lab_tb = s.shapes.add_textbox(left, top, cw, Inches(0.3))
        _add_par(lab_tb.text_frame, lab, size=12, bold=True, color=(ACCENT if i == 3 else NAVY), first=True)
        p = os.path.join(fe, fn)
        if os.path.exists(p):
            place_image(s, p, left, top + Inches(0.32), cw, ch - Inches(0.1))
    caption(s, "Gold = 4, model said 3. It counted the three store displays (Videos 1–3) but ignored the "
               "drinks inside the cart (Video 4). These are the 8 frames the model actually saw per clip.",
            top=SH - Inches(0.7))


def s_modes():
    bullets_slide("Named failure modes (verbatim from reasoning traces)", [
        "**Wrong cross-view aggregation** — counting/find-the-common-item across clips; misses the non-obvious shared object (case above)",
        "**Single-clip shortcutting** — answers from one informative view, ignores the rest",
        "**Given-order bias** — \"put the events in order\" defaults to clip order 1-2-3-4 instead of reasoning it out",
        "**Yes-bias / premise acceptance** — agrees with what the question assumes without checking the video",
        "**Temporal grounding** — confuses video-level timestamps with event-level time (the shared weak spot)",
        "Full traces: regenerate via analyze_failures.py (qwen3vl/internvl3_failures.md)",
    ], body_size=16)


def s_datasets():
    bullets_slide("Datasets — where the camera views come from", [
        "✓ **MEVA** — ready; public S3, no credentials. The 60-Q + 1,033-Q pools; 100% answer-safe under the <=4 cap",
        "✓ **EgoExo4D** — license approved & videos NOW DOWNLOADED (96/96 files local, readable); combined 100-Q eval staged, not yet run",
        "**AgiBot** — wired; gated behind a HF click-through + ~810 GB of tars; deferred",
        "**nuScenes** — skipped (fixed 6 cameras → always lossy under the <=4 cap)",
        "Key idea: questions can need up to 16 cameras but the harness feeds <=4 → we tag cap_answer_safe and keep the true camera count (orig_num_cameras) for the curve",
        "Walkthroughs: egoexo_setup.md · agibot_setup.md",
    ], body_size=15)


def s_status_full():
    bullets_slide("Status & next steps", [
        "**Done this round:**",
        ("clean Qwen re-run (scoring bug fixed) · scaled CrossView to 1,033 Q on both models · failure taxonomy from the traces · CrossView blind baseline · EgoExo4D downloaded · full input/output pipeline documented (ORIENTATION.md)", 1),
        "**Next:**",
        ("run the combined MEVA + EgoExo4D eval (videos are local) — first cross-source comparison", 1),
        ("equalize the two harnesses (same frame sampling + stop settings) for a truly fair head-to-head", 1),
        ("accuracy-vs-token-budget / #frames curves anchored on the blind 0-token floor", 1),
    ])


def s_status_concise():
    bullets_slide("Status & next steps", [
        "**Done this round:** clean Qwen re-run (scoring bug fixed) · CrossView scaled to **1,033 Q** on both models · failure taxonomy · CrossView blind baseline · **EgoExo4D downloaded**",
        "**Next:** run the combined MEVA + EgoExo4D eval (videos are local); equalize the two harnesses for a fair head-to-head",
        "**Where to look:** analysis/ORIENTATION.md (full guide) · crossview_meva1033_out/failure_examples_curated.md",
    ], body_size=17)


def s_wheretolook():
    bullets_slide("Where to look", [
        "**analysis/ORIENTATION.md** — the single onboarding guide: what we did, how inputs & outputs are enabled (with file:line anchors), how to run it, results inventory, findings",
        "analysis/README.md — pipeline index + headline numbers",
        "analysis/mentor_report.md — failure-mode analysis (the deliverable)",
        "analysis/crossview_meva1033_out/failure_examples_curated.md — the failure taxonomy digest",
        "Qwen harness: Video-R1/src/eval_thinking.py · InternVL: lmms-eval crossview_think task",
    ], body_size=16)


FULL = [s_title, s_what_full, s_benchmarks, s_inputs, s_outputs, s_headline, s_flip,
        s_eofix, s_cameras, s_taskflip, s_mech_full, s_case, s_modes, s_datasets,
        s_status_full, s_wheretolook]
CONCISE_DECK = [s_title, s_what_concise, s_headline, s_flip, s_eofix, s_cameras,
                s_mech_concise, s_status_concise]

for build in (CONCISE_DECK if CONCISE else FULL):
    build()

prs.save(OUT)
print("wrote", OUT, "·", len(prs.slides._sldIdLst), "slides", "(concise)" if CONCISE else "(full)")

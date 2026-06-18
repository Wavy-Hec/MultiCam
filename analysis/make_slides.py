#!/usr/bin/env python3
"""Build the multi-camera VLM progress deck (python-pptx) for mentor feedback.

Plain-language, ~5 bullets/slide, jargon glossed on first use. Directly answers
"can models answer without reasoning?" (slides 4-5) and is honest about the Qwen
CrossView scoring artifact (slide 8) — CrossView charts show InternVL3 only.

Run:  ~/anaconda3/envs/cvbench/bin/python analysis/make_slides.py
Out:  analysis/multicam_progress.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
A = os.path.join(REPO, "analysis")
OUT = os.path.join(A, "multicam_progress.pptx")

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1E, 0x7E, 0x34)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _add_par(tf, text, size=18, bold=False, color=GRAY, level=0, first=False, space_after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size, run.font.bold, run.font.color.rgb = Pt(size), bold, color
    return p


def _title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), SW - Inches(1.2), Inches(1.0))
    _add_par(tb.text_frame, text, size=29, bold=True, color=NAVY, first=True)


def _fill_bullets(tf, items, body_size):
    first = True
    for it in items:
        text, level = it if isinstance(it, tuple) else (it, 0)
        bold = text.startswith("**")
        text = text.replace("**", "")
        color = ACCENT if text.startswith("⚠") else (GREEN if text.startswith("✓") else GRAY)
        prefix = "" if (bold and level == 0) else ("•  " if level == 0 else "–  ")
        _add_par(tf, prefix + text, size=body_size - 2 * level, bold=bold,
                 color=color, level=level, first=first, space_after=8)
        first = False


def bullets_slide(title, items, body_size=18):
    s = prs.slides.add_slide(BLANK)
    _title(s, title)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.55), SW - Inches(1.4), SH - Inches(2.0))
    tb.text_frame.word_wrap = True
    _fill_bullets(tb.text_frame, items, body_size)
    return s


def place_image(slide, path, left, top, max_w, max_h):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w, h = max_w, max_w / ar
    if h > max_h:
        h, w = max_h, max_h * ar
    slide.shapes.add_picture(path, Emu(int(left + (max_w - w) / 2)),
                             Emu(int(top + (max_h - h) / 2)), width=Emu(int(w)), height=Emu(int(h)))


def caption(slide, text, color=GRAY, top=None):
    top = top if top is not None else SH - Inches(1.0)
    tb = slide.shapes.add_textbox(Inches(0.7), top, SW - Inches(1.4), Inches(0.85))
    tb.text_frame.word_wrap = True
    _add_par(tb.text_frame, text, size=15, color=color, first=True)


def image_slide(title, img, cap, cap_color=GRAY):
    s = prs.slides.add_slide(BLANK)
    _title(s, title)
    place_image(s, img, Inches(0.7), Inches(1.45), SW - Inches(1.4), Inches(4.55))
    caption(s, cap, color=cap_color)
    return s


def image_then_bullets(title, img, bullets, body_size=16, img_h=Inches(3.95)):
    s = prs.slides.add_slide(BLANK)
    _title(s, title)
    place_image(s, img, Inches(0.7), Inches(1.4), SW - Inches(1.4), img_h)
    btop = Inches(1.4) + img_h + Inches(0.15)
    tb = s.shapes.add_textbox(Inches(0.7), btop, SW - Inches(1.4), SH - btop - Inches(0.2))
    tb.text_frame.word_wrap = True
    _fill_bullets(tb.text_frame, bullets, body_size)
    return s


def title_slide(title, sub_lines):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.4), SW - Inches(1.6), Inches(2.0))
    tb.text_frame.word_wrap = True
    _add_par(tb.text_frame, title, size=38, bold=True, color=NAVY, first=True)
    for ln in sub_lines:
        _add_par(tb.text_frame, ln, size=20, color=GRAY, space_after=4)
    return s


# ----------------------------------------------------------------- slides
title_slide("Multi-Camera VLM Benchmarking — Progress Update",
            ["CVBench  +  CrossView (UT Austin multi-camera dataset)",
             "Hector Lugo  ·  for Dr. Datta & Harsh — feedback welcome",
             "2026-06-16"])

bullets_slide("What we're testing", [
    "Two \"thinking\" vision-language models — they write out their reasoning before answering, so we can read WHY they fail",
    ("Qwen3-VL-8B  and  InternVL3-8B", 1),
    "A \"blind\" control: the same questions with the video removed (text only)",
    "Two benchmarks, side by side: CVBench (unrelated clips) and CrossView (same scene, many synchronized cameras)",
    "Goal: not just a score — an interpretable picture of how the models fail",
])

bullets_slide("Two benchmarks (different sizes), run side by side", [
    "**CVBench** — questions over 1–4 UNRELATED video clips   (we ran 45 questions)",
    "**CrossView** — questions over the SAME scene from many synchronized cameras, from UT Austin   (we ran 60 questions)",
    "CVBench asks \"can it connect separate clips?\";  CrossView asks \"can it fuse simultaneous camera views?\"",
    "Both run through the identical harness → results are directly comparable",
])

bullets_slide("Does 40% mean the models answer without reasoning?", [
    "**Short answer: partly — and where it happens, it's the question wording, not understanding the video.**",
    "Text-only (no video) scores 40%.  Random guessing on these questions ≈ 28–32%.",
    "So text-only beats guessing by only ~9 points — and that edge comes from wording:",
    ("a few question types have give-away answer options", 1),
    ("\"yes-bias\": the model tends to agree with whatever the question assumes", 1),
    "The blind test is exactly how we DETECT those guessable questions, so we can flag / remove them",
    "Proof the models DO use the video where it matters  →  next slide",
])

image_then_bullets("How much the answer needs the video", os.path.join(A, "reasoning_gain.png"), [
    "**2 clips: video adds only +7 pts** (easy, often yes/no — the guessable set)",
    "**3–4 clips: video adds ~+30 pts** — text-only falls toward chance; the video carries the answer",
    "The gap proves the answer NEEDS the video; the reasoning traces (failure slides) show how it reasons over it",
    "Caveat: small samples (~15/group), single run → the JUMP at 3+ clips is the robust trend, not the exact points",
])

bullets_slide("The blind (text-only) control is clean", [
    "✓ No image data of any kind reached the model — checked on all 45 prompts",
    "✓ No hidden leaks: no filenames, IDs, timestamps, or frame hints in the text",
    "✓ The 40% was recounted three independent ways — the number holds up",
    "Caveat: the questions still mention videos in the options, so the model knows videos EXIST — it just gets none of their content",
    "Detail: analysis/blind_sanity.md",
])

bullets_slide("How the video reaches the model", [
    "Each clip is reduced to 8 still frames, sampled evenly across the clip",
    "Clips are shown in order, with a separator marking where each one starts:",
    ("Qwen: a text label  (\"Video 1:\", \"Video 2:\" …)", 1),
    ("InternVL: a small marker image that says  \"This is video 1\"", 1),
    "So the model never sees full video — just 8 frames per clip, in order, then the question",
    "Detail: analysis/input_pipeline.md",
])

bullets_slide("⚠ Heads-up: Qwen's CrossView scores are a scoring bug, not a result", [
    "**Bottom line: ignore Qwen's CrossView numbers for now — the next charts show InternVL3 only.**",
    "What happened: on the ordering questions, Qwen's answer ran past the max output length and never finished",
    "With no final answer to read, the scorer defaulted to \"A\" — but the correct answer is never \"A\", so it scored 0",
    "The same default also inflated Qwen's other CrossView scores → unreliable across the board",
    "✓ InternVL3 finished and was scored correctly → its numbers ARE trustworthy",
    "Fix in progress: longer output limit + a scorer that abstains instead of guessing, then re-run",
])

image_slide("CrossView: accuracy drops as more cameras are involved",
            os.path.join(A, "crossview_out", "accuracy_vs_orig_cameras_internvl3.png"),
            "InternVL3: 65% at 2 cameras → 0% at 5. Direct evidence that fusing many simultaneous "
            "views is hard. (Noisy at high camera counts — small samples. Qwen excluded — see previous slide.)")

image_slide("CrossView: accuracy by question type (InternVL3)",
            os.path.join(A, "crossview_out", "accuracy_by_task_internvl3.png"),
            "Weakest on Event-Ordering (30%) and Spatial (40%); Temporal 45%. 20 questions per type.")

# ---- failure case with input frames (2x2 grid) ----
s = prs.slides.add_slide(BLANK)
_title(s, "Failure case: counting across 4 cameras")
fe = os.path.join(A, "failure_examples")
cells = [("Video 1 — store display", "id57_video1_frames.jpg"),
         ("Video 2 — store display", "id57_video2_frames.jpg"),
         ("Video 3 — store display", "id57_video3_frames.jpg"),
         ("Video 4 — drinks INSIDE the cart (missed)", "id57_video4_frames.jpg")]
gx, gy = Inches(0.7), Inches(1.45)
cw, ch = (SW - Inches(1.4)) / 2, Inches(2.0)
for i, (lab, fn) in enumerate(cells):
    r, c = divmod(i, 2)
    left, top = gx + c * cw, gy + r * (ch + Inches(0.32))
    lab_tb = s.shapes.add_textbox(left, top, cw, Inches(0.3))
    _add_par(lab_tb.text_frame, lab, size=12, bold=True, color=(ACCENT if i == 3 else NAVY), first=True)
    p = os.path.join(fe, fn)
    if os.path.exists(p):
        place_image(s, p, left, top + Inches(0.32), cw, ch - Inches(0.1))
caption(s, "Gold = 4, model said 3. It counted the three store displays (Videos 1–3) but ignored the "
           "drinks inside the cart (Video 4). These are the 8 frames the model actually saw per clip.",
        top=SH - Inches(0.7))

bullets_slide("Common failure modes (read from the reasoning traces)", [
    "**Agrees with the question's assumption** (\"yes-bias\") — answers without really checking the video",
    "**Looks at only one clip** and ignores the others",
    "**Counting across clips** — misses the shared item that isn't the obvious subject (the case above)",
    "**\"Put the events in order\"** — defaults to the clip order 1-2-3-4 instead of working it out",
    "Full traces: cvbench_multicam_failures.md · crossview_out/internvl3_failures.md",
])

bullets_slide("Status & next steps", [
    "**Done:**",
    ("both benchmarks run · the 3 analysis tasks · failure catalog · charts · caught the Qwen scoring bug", 1),
    "**In progress:**",
    ("Qwen re-run (cluster delays) → will restore Qwen to the CrossView charts", 1),
    "**Next:**",
    ("add a blind control for CrossView (is it harder to guess than CVBench?)", 1),
    ("bring in more camera sources beyond MEVA; larger samples for solid per-type numbers", 1),
])

prs.save(OUT)
print("wrote", OUT, "·", len(prs.slides._sldIdLst), "slides")
